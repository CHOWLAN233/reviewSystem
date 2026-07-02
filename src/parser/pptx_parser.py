"""
PPTX / PPT file parser.

Extracts text from:
- .pptx: python-pptx (Open XML format)
- .ppt : OLE compound document binary extraction (legacy format)

Also extracts speaker notes from .pptx files.
"""

from __future__ import annotations

import logging
import re
import struct
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table

from .base_parser import AbstractParser

logger = logging.getLogger(__name__)


class PPTXParserError(Exception):
    """Raised when a PPT/PPTX file cannot be parsed."""


class PPTXParser(AbstractParser):
    """Extract text from PowerPoint (.pptx and legacy .ppt) files."""

    def extract_text(self, filepath: Path) -> str:
        """Try python-pptx first; fall back to binary extraction for legacy .ppt."""
        try:
            return self._extract_text_pptx(filepath)
        except PPTXParserError:
            if filepath.suffix.lower() == ".ppt":
                logger.info(f"Legacy .ppt format detected, using binary extraction: {filepath.name}")
                return self._extract_text_legacy_ppt(filepath)
            raise

    def extract_pages_text(self, filepath: Path, start: int, count: int) -> str:
        """Extract text from slides/pages [*start*, *start* + *count*)."""
        try:
            return self._extract_pages_pptx(filepath, start, count)
        except PPTXParserError:
            if filepath.suffix.lower() == ".ppt":
                logger.info(f"Legacy .ppt format, using binary extraction for first {count} slides")
                # Binary extraction doesn't have slide-level granularity;
                # return a prefix of the full text as an approximation.
                full = self._extract_text_legacy_ppt(filepath)
                # Split on slide-like boundaries (common in extracted PPT text)
                slides = re.split(r"\n*---\s*Slide\s*\d+\s*---\n*", full)
                # First element is usually empty or pre-content
                if slides and not slides[0].strip():
                    slides = slides[1:]
                selected = slides[start:start + count]
                return "\n\n".join(selected) if selected else full[:3000]
            raise

    def get_page_count(self, filepath: Path) -> int:
        """Return the number of slides (estimated for legacy .ppt)."""
        try:
            prs = Presentation(str(filepath))
            return len(prs.slides)
        except Exception:
            # For legacy .ppt, estimate from slide markers in extracted text
            text = self._extract_text_legacy_ppt(filepath)
            slides = re.findall(r"---\s*Slide\s*\d+\s*---", text)
            return max(len(slides), 1)

    # ------------------------------------------------------------------
    # .pptx (Open XML) extraction
    # ------------------------------------------------------------------

    def _extract_text_pptx(self, filepath: Path) -> str:
        """Extract text from a .pptx file using python-pptx."""
        prs = self._open_pptx(filepath)
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = [f"\n--- Slide {idx} ---"]
            for shape in slide.shapes:
                text = self._extract_shape_text(shape)
                if text.strip():
                    slide_lines.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"[Speaker Notes]: {notes}")
            parts.append("\n".join(slide_lines))
        return "\n".join(parts)

    def _extract_pages_pptx(self, filepath: Path, start: int, count: int) -> str:
        """Extract text from .pptx slides [*start*, *start* + *count*)."""
        prs = self._open_pptx(filepath)
        total = len(prs.slides)
        end = min(start + count, total)
        parts: list[str] = []
        for idx in range(start, end):
            slide = prs.slides[idx]
            slide_lines: list[str] = [f"\n--- Slide {idx + 1} ---"]
            for shape in slide.shapes:
                text = self._extract_shape_text(shape)
                if text.strip():
                    slide_lines.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"[Speaker Notes]: {notes}")
            parts.append("\n".join(slide_lines))
        return "\n".join(parts)

    def _open_pptx(self, filepath: Path) -> Presentation:
        """Open a .pptx file, raising PPTXParserError on failure."""
        try:
            return Presentation(str(filepath))
        except Exception as exc:
            raise PPTXParserError(f"Failed to open {filepath.name}: {exc}") from exc

    # ------------------------------------------------------------------
    # .ppt (legacy OLE binary) extraction
    # ------------------------------------------------------------------

    # Maximum characters to extract from legacy .ppt files
    # Increased to match the summarizer's larger chunk size (MAX_CHUNK_CHARS=200K).
    # Modern LLMs handle much more text, so there's no need to truncate aggressively.
    MAX_LEGACY_CHARS = 400000

    def _extract_text_legacy_ppt(self, filepath: Path) -> str:
        """
        Extract readable text from a legacy .ppt (OLE compound document).

        Strategy:
        1. Try ``olefile`` package (robust OLE parsing + text record extraction).
        2. Fall back to raw binary string extraction (extract ASCII/UTF-8 sequences).
        """
        # Strategy 1: olefile
        try:
            import olefile
            return self._extract_ppt_via_olefile(filepath)
        except ImportError:
            logger.debug("olefile not installed, using raw binary extraction.")
        except Exception as exc:
            logger.debug(f"olefile extraction failed: {exc}")

        # Strategy 2: raw binary string scraping
        return self._extract_ppt_via_raw(filepath)

    def _extract_ppt_via_olefile(self, filepath: Path) -> str:
        """Extract text from .ppt using olefile for structured OLE reading."""
        import olefile

        ole = olefile.OleFileIO(str(filepath))
        parts: list[str] = []

        # PowerPoint stores text in records within "PowerPoint Document" stream
        if ole.exists("PowerPoint Document"):
            data = ole.openstream("PowerPoint Document").read()
            # Extract UTF-16LE text records (PowerPoint uses UTF-16LE internally)
            text_records = self._extract_utf16le_records(data)
            if text_records:
                for i, record in enumerate(text_records, 1):
                    if record.strip():
                        parts.append(f"\n--- Slide {i} ---\n{record}")
            else:
                # Fall back to extracting any readable text
                parts.append(self._extract_readable_strings(data))

        # Also check for "Current User" stream's text
        for stream_name in ole.listdir():
            name = "/".join(stream_name) if isinstance(stream_name, (list, tuple)) else str(stream_name)
            if "PowerPoint Document" not in name:
                continue

        ole.close()

        result = "\n".join(parts) if parts else self._extract_ppt_via_raw(filepath)
        return self._truncate_text(result)

    def _extract_ppt_via_raw(self, filepath: Path) -> str:
        """Extract readable text from .ppt binary by scraping for strings."""
        data = filepath.read_bytes()

        # Try UTF-16LE text extraction (PowerPoint's native encoding)
        utf16_text = self._extract_utf16le_records(data)
        if utf16_text:
            parts = []
            slide_num = 1
            for record in utf16_text:
                record = record.strip()
                if record and len(record) > 3:
                    parts.append(f"\n--- Slide {slide_num} ---\n{record}")
                    slide_num += 1
            if parts:
                return "\n".join(parts)

        # Fall back to ASCII/UTF-8 string extraction
        text = self._extract_readable_strings(data)
        return self._truncate_text(text) if text.strip() else "(No readable text extracted from legacy .ppt file)"

    def _truncate_text(self, text: str) -> str:
        """Truncate text to MAX_LEGACY_CHARS while preserving slide boundaries."""
        if len(text) <= self.MAX_LEGACY_CHARS:
            return text
        # Keep the beginning and note the truncation
        truncated = text[:self.MAX_LEGACY_CHARS]
        last_slide = truncated.rfind("\n--- Slide")
        if last_slide > 0:
            truncated = truncated[:last_slide]
        logger.warning(
            f"Truncated legacy PPT text from {len(text):,} to {len(truncated):,} chars"
        )
        return truncated + "\n\n[... Content truncated due to length ...]"

    @staticmethod
    def _extract_utf16le_records(data: bytes, min_len: int = 4) -> list[str]:
        """
        Extract UTF-16LE text records from binary data.

        PowerPoint stores text as UTF-16LE encoded strings preceded by
        a 4-byte record header. We scan for sequences of valid UTF-16LE
        characters and collect them.
        """
        records: list[str] = []
        i = 0
        current: list[str] = []
        consecutive_null = 0

        while i < len(data) - 1:
            # Read a UTF-16LE character (2 bytes)
            char_code = struct.unpack_from("<H", data, i)[0]

            if char_code == 0:
                consecutive_null += 1
                # Two consecutive nulls (4 zero bytes) = likely record boundary
                if consecutive_null >= 2 and current:
                    text = "".join(current).strip()
                    if len(text) >= min_len:
                        records.append(text)
                    current = []
                    consecutive_null = 0
            elif 0x20 <= char_code <= 0xFFEF or char_code in (0x0D, 0x0A, 0x09):
                # Printable Unicode range (including CJK, Latin, etc.)
                try:
                    char = chr(char_code)
                    if char.isprintable() or char in "\r\n\t":
                        current.append(char)
                        consecutive_null = 0
                    else:
                        consecutive_null += 1
                except ValueError:
                    consecutive_null += 1
            else:
                consecutive_null += 1

            i += 2

        # Don't forget the last record
        if current:
            text = "".join(current).strip()
            if len(text) >= min_len:
                records.append(text)

        return records

    @staticmethod
    def _extract_readable_strings(data: bytes, min_len: int = 4) -> str:
        """
        Extract readable text sequences from binary data.

        Collects consecutive printable ASCII/CJK characters.
        """
        result: list[str] = []
        current: list[str] = []

        for byte in data:
            char = chr(byte)
            if char.isprintable() or char in "\r\n\t ":
                current.append(char)
            else:
                if len(current) >= min_len:
                    text = "".join(current).strip()
                    if text:
                        result.append(text)
                current = []

        if len(current) >= min_len:
            text = "".join(current).strip()
            if text:
                result.append(text)

        return "\n".join(result)

    # ------------------------------------------------------------------
    # Shape text extraction (for .pptx)
    # ------------------------------------------------------------------

    def _extract_shape_text(self, shape: BaseShape) -> str:
        """Extract text from a shape, handling text frames, tables, and groups."""
        # Text frame (most common)
        if shape.has_text_frame:
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                p_text = para.text.strip()
                if p_text:
                    paragraphs.append(p_text)
            return "\n".join(paragraphs)

        # Table
        if shape.has_table:
            table: Table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(c for c in cells if c))
            return "\n".join(rows)

        # Group shape – recurse
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                return "\n".join(
                    self._extract_shape_text(s) for s in shape.shapes
                )
            except AttributeError:
                pass

        return ""
